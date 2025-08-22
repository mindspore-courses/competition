import ray
import re
import PyPDF2
import markdown
from bs4 import BeautifulSoup
import tiktoken
from typing import List

from io import BytesIO
import logging

from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter

import docx
import pptx
import pandas as pd
import whisper
import soundfile as sf
import tempfile
import os

# ==============================================================================
# 1. å®šä¹‰ Embedding Actor
# (é€»è¾‘æ¥è‡ªåŸ embedding_server/app.py)
# ==============================================================================
@ray.remote(num_cpus=4)  # ä¸ºæ¯ä¸ª Embedding Actor å®ä¾‹åˆ†é…4ä¸ªCPU
class EmbeddingActor:
    """
    ä¸€ä¸ªä¸“ç”¨äºæ–‡æœ¬å‘é‡åŒ–çš„ Ray Actorã€‚
    å®ƒåœ¨è‡ªå·±çš„è¿›ç¨‹ä¸­åŠ è½½å¹¶æŒæœ‰ BAAI/bge-base-zh-v1.5 æ¨¡å‹ã€‚
    """
    def __init__(self):
        # åœ¨ Actor åˆå§‹åŒ–æ—¶åŠ è½½æ¨¡å‹ï¼Œæ¨¡å‹å°†å¸¸é©»äºè¯¥ Actor çš„æ˜¾å­˜ä¸­
        try:
            from mindnlp.sentence import SentenceTransformer
            print("â–¶ï¸ EmbeddingActor: æ­£åœ¨åŠ è½½æ¨¡å‹ (BAAI/bge-base-zh-v1.5)...")
            self.model = SentenceTransformer('BAAI/bge-base-zh-v1.5')
            print("âœ… EmbeddingActor: æ¨¡å‹åŠ è½½æˆåŠŸ!")
        except Exception as e:
            print(f"âŒ EmbeddingActor: åŠ è½½æ¨¡å‹å¤±è´¥: {e}")
            self.model = None

    def embed(self, texts: List[str]) -> List[List[float]]:
        """æ¥æ”¶æ–‡æœ¬åˆ—è¡¨ï¼Œè¿”å›å‘é‡åˆ—è¡¨ã€‚è¿™æ˜¯è¯¥ Actor çš„æ ¸å¿ƒæ¨ç†æ–¹æ³•ã€‚"""
        if not self.model or not texts:
            return []
        
        print(f"âš™ï¸ EmbeddingActor: æ­£åœ¨ä¸º {len(texts)} æ¡æ–‡æœ¬ç”Ÿæˆå‘é‡...")
        try:
            vectors = self.model.encode(texts, normalize_embeddings=True)
            return vectors.tolist()
        except Exception as e:
            print(f"âŒ EmbeddingActor: å‘é‡åŒ–è¿‡ç¨‹ä¸­å‘ç”Ÿé”™è¯¯: {e}")
            return []

# ==============================================================================
# 2. å®šä¹‰ LLM Actor
# (é€»è¾‘æ¥è‡ªåŸ llm_server/app.py)
# ==============================================================================
@ray.remote(num_cpus=6) # ä¸ºæ¯ä¸ª LLM Actor å®ä¾‹åˆ†é…6ä¸ªCPU
class LLMActor:
    """
    ä¸€ä¸ªä¸“ç”¨äºå¤§è¯­è¨€æ¨¡å‹æ¨ç†çš„ Ray Actorã€‚
    å®ƒåœ¨è‡ªå·±çš„è¿›ç¨‹ä¸­åŠ è½½å¹¶æŒæœ‰ MiniCPM æ¨¡å‹ã€‚
    """
    def __init__(self):
        # åœ¨ Actor åˆå§‹åŒ–æ—¶åŠ è½½æ¨¡å‹å’Œåˆ†è¯å™¨
        try:
            import mindspore
            from mindnlp.transformers import AutoTokenizer, AutoModelForCausalLM
            
            LLM_MODEL_PATH = 'openbmb/MiniCPM-2B-dpo-bf16'
            logging.info(f"â–¶ï¸ LLMActor: æ­£åœ¨åŠ è½½LLMæ¨¡å‹åŠåˆ†è¯å™¨ ({LLM_MODEL_PATH})...")
            
            self.tokenizer = AutoTokenizer.from_pretrained(LLM_MODEL_PATH, mirror="huggingface")
            self.model = AutoModelForCausalLM.from_pretrained(LLM_MODEL_PATH, ms_dtype=mindspore.float32, mirror="huggingface")
            
            logging.info("âœ… LLMActor: æ¨¡å‹åŠ è½½æˆåŠŸ!")
        except Exception as e:
            logging.info(f"âŒ LLMActor: åŠ è½½LLMæ¨¡å‹æ—¶å‘ç”Ÿä¸¥é‡é”™è¯¯: {e}")
            self.model = None
            self.tokenizer = None

    def generate(self, prompt: str, max_length: int = 1024) -> str:
        """æ¥æ”¶æ ¼å¼åŒ–åçš„ promptï¼Œè¿”å›ç”Ÿæˆçš„æ–‡æœ¬ã€‚"""
        if not self.model or not self.tokenizer:
            return "LLM Actor æ¨¡å‹æœªæˆåŠŸåŠ è½½ï¼Œæ— æ³•ç”Ÿæˆå›ç­”ã€‚"
            
        print(f"âš™ï¸ LLMActor: æ”¶åˆ°ç”Ÿæˆè¯·æ±‚ï¼Œæ­£åœ¨è°ƒç”¨ model.chat...")
        try:
            response_text, _ = self.model.chat(
                self.tokenizer,
                prompt,
                history=[],
                max_length=max_length,
            )
            return response_text
        except Exception as e:
            print(f"âŒ LLMActor: æ¨ç†è¿‡ç¨‹ä¸­å‘ç”Ÿé”™è¯¯: {e}")
            return f"LLM Actor æ¨ç†å¤±è´¥: {str(e)}"

# ==============================================================================
# 3. å®šä¹‰æ–‡ä»¶å¤„ç† Task
# (é€»è¾‘æ¥è‡ªåŸ main_app/main.py ä¸­çš„ FileProcessor ç±»)
# è¿™æ˜¯ä¸€ä¸ªæ— çŠ¶æ€çš„ä»»åŠ¡ï¼Œéå¸¸é€‚åˆç”¨ Ray Task æ¥å¹¶è¡Œå¤„ç†ã€‚
# ==============================================================================
@ray.remote
def parse_and_chunk_document(file_content: bytes, file_name: str) -> List[str]:
    """
    ä¸€ä¸ª Ray Taskï¼Œç”¨äºè§£æå•ä¸ªæ–‡ä»¶å†…å®¹å¹¶å°†å…¶åˆ†å—ã€‚
    æ­¤ç‰ˆæœ¬é‡‡ç”¨äº†æ›´æ™ºèƒ½çš„ã€æ„ŸçŸ¥å†…å®¹ç»“æ„çš„åˆ†å—ç­–ç•¥ã€‚
    """
    print(f"âš™ï¸ Ray Task: æ­£åœ¨è§£ææ–‡ä»¶ '{file_name}'...")
    from rapidocr_onnxruntime import RapidOCR
    from PIL import Image
    import numpy as np
    ocr_engine = RapidOCR()

    def read_pdf(stream) -> str:
        reader = PyPDF2.PdfReader(stream)
        return "".join(page.extract_text() for page in reader.pages if page.extract_text())

    def read_markdown(stream) -> str:
        md_text = stream.read().decode('utf-8')
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, 'html.parser')
        return re.sub(r'http\S+', '', soup.get_text())

    def read_text(stream) -> str:
        return stream.read().decode('utf-8')

    def read_image(stream) -> str:
        from rapidocr_onnxruntime import RapidOCR
        from PIL import Image
        import numpy as np
        ocr_engine = RapidOCR()
        img = Image.open(stream)
        img_np = np.array(img)
        result, _ = ocr_engine(img_np)
        if result:
            return "\n".join([line[1] for line in result])
        return ""

    def read_docx(stream) -> str:
        document = docx.Document(stream)
        return "\n".join([para.text for para in document.paragraphs])

    def read_pptx(stream) -> str:
        presentation = pptx.Presentation(stream)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs)

    def read_csv(stream) -> str:
        df = pd.read_csv(stream)
        return df.to_string()

    def read_audio(stream, file_name) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file_name)[1]) as tmp_file:
            tmp_file.write(stream.read())
            tmp_file_path = tmp_file.name
        
        try:
            print("ğŸ¤ Ray Task: æ­£åœ¨åŠ è½½ Whisper æ¨¡å‹è¿›è¡Œè¯­éŸ³è¯†åˆ«...")
            model = whisper.load_model("small")
            result = model.transcribe(tmp_file_path, fp16=False)
            print(f"ğŸ”Š Ray Task: æ–‡ä»¶ '{file_name}' è¯­éŸ³è¯†åˆ«å®Œæˆã€‚")
            return result.get("text", "")
        finally:
            os.remove(tmp_file_path) 


    file_stream = BytesIO(file_content)
    text = ""
    file_suffix = file_name.split('.')[-1].lower()

    if file_suffix == 'pdf': text = read_pdf(file_stream)
    elif file_suffix == 'md': text = read_markdown(file_stream)
    elif file_suffix == 'txt': text = read_text(file_stream)
    elif file_suffix in ['png', 'jpg', 'jpeg']: text = read_image(file_stream)
    elif file_suffix == 'docx': text = read_docx(file_stream)
    elif file_suffix == 'pptx': text = read_pptx(file_stream)
    elif file_suffix == 'csv': text = read_csv(file_stream)
    elif file_suffix in ['wav', 'mp3', 'm4a']: text = read_audio(file_stream, file_name)
    else:
        print(f"âš ï¸ Ray Task: ä¸æ”¯æŒçš„æ–‡ä»¶ç±»å‹ '{file_suffix}'ï¼Œè·³è¿‡æ–‡ä»¶ {file_name}ã€‚")
        return []

    if not text:
        print(f"â„¹ï¸ Ray Task: ä»æ–‡ä»¶ '{file_name}' ä¸­æœªæå–åˆ°æ–‡æœ¬ã€‚")
        return []

    chunks = []

    if file_suffix == 'md':
        # å¯¹ Markdown æ–‡ä»¶ä½¿ç”¨æ ‡é¢˜åˆ†å‰²å™¨
        print(f"âœ¨ Ray Task: å¯¹ Markdown æ–‡ä»¶ '{file_name}' ä½¿ç”¨æ ‡é¢˜åˆ†å‰²ç­–ç•¥ã€‚")
        headers_to_split_on = [
            ("#", "Header 1"),
            ("##", "Header 2"),
            ("###", "Header 3"),
        ]
        markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on, strip_headers=False)
        md_header_splits = markdown_splitter.split_text(text)
        
        # å¯¹åˆ†å‰²åçš„æ¯ä¸ªå¤§å—ï¼Œå†è¿›è¡Œé€’å½’åˆ†å—ï¼Œé˜²æ­¢æœ‰è¶…é•¿ç« èŠ‚
        # (è¿™éƒ¨åˆ†å’Œä¸‹é¢çš„é€’å½’åˆ†å—é€»è¾‘æ˜¯å¤ç”¨çš„)
        chunk_size = 600
        chunk_overlap = 150
        text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            encoding_name="cl100k_base",
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
        for split in md_header_splits:
            chunks.extend(text_splitter.split_text(split.page_content))

    else:
        # å¯¹å…¶ä»–ç±»å‹æ–‡ä»¶ï¼ˆPDF, TXT, å›¾ç‰‡OCRç»“æœï¼‰ä½¿ç”¨é€’å½’å­—ç¬¦åˆ†å‰²å™¨
        print(f"âœ¨ Ray Task: å¯¹æ–‡ä»¶ '{file_name}' ä½¿ç”¨é€’å½’å­—ç¬¦åˆ†å‰²ç­–ç•¥ã€‚")
        chunk_size = 600
        chunk_overlap = 150
        # ä½¿ç”¨tiktokenæ¥è®¡ç®—å—é•¿åº¦
        text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            encoding_name="cl100k_base",
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
        chunks = text_splitter.split_text(text)
        
    final_chunks = [
        {
            "content": chunk_text,
            "source": file_name,
            "chunk_index": i + 1
        }
        for i, chunk_text in enumerate(chunks)
    ]
        
    print(f"âœ… Ray Task: æ–‡ä»¶ '{file_name}' è§£æå¹¶åˆ†å—ä¸º {len(final_chunks)} å—ã€‚")
    return final_chunks