import ray
import re
import PyPDF2
import markdown
from bs4 import BeautifulSoup
import tiktoken
from typing import List

from io import BytesIO
import logging

from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter

import docx
import pptx
import pandas as pd
import whisper
import soundfile as sf
import tempfile
import os

# 1.  Embedding Actor
@ray.remote(num_cpus=4) 
class EmbeddingActor:
    def __init__(self):
        try:
            from mindnlp.sentence import SentenceTransformer
            print("▶️ EmbeddingActor: 正在加载模型 (BAAI/bge-base-zh-v1.5)...")
            MODEL_PATH = '/app/.mindnlp/model/BAAI/bge-base-zh-v1.5'
            MODEL_PATH = 'BAAI/bge-base-zh-v1.5' # 自动下载到容器内：/app/.mindnlp/model -> 挂载到DistributedRAG/rag_models_cache
            self.model = SentenceTransformer(MODEL_PATH)
            print("✅ EmbeddingActor: 模型加载成功!")
            print(self.model)   
        except Exception as e:
            print(f"❌ EmbeddingActor: 加载模型失败: {e}")
            self.model = None

    def embed(self, texts: List[str]) -> List[List[float]]:
        if not self.model or not texts:
            return []
        
        print(f"⚙️ EmbeddingActor: 正在为 {len(texts)} 条文本生成向量...")
        try:
            vectors = self.model.encode(texts, normalize_embeddings=True)
            return vectors.tolist()
        except Exception as e:
            print(f"❌ EmbeddingActor: 向量化过程中发生错误: {e}")
            return []

# 2. LLM Actor
@ray.remote(num_cpus=6)
class LLMActor:
    def __init__(self):
        try:
            import mindspore
            from mindnlp.transformers import AutoTokenizer, AutoModelForCausalLM
            LLM_MODEL_PATH = 'openbmb/MiniCPM-2B-dpo-bf16'
            # LLM_MODEL_PATH = '/app/.mindnlp/model/openbmb/MiniCPM-2B-dpo-bf16' # 自动下载到容器内：/app/.mindnlp/model -> 挂载到DistributedRAG/rag_models_cache
            logging.info(f"▶️ LLMActor: 正在加载LLM模型及分词器 ({LLM_MODEL_PATH})...")
            
            self.tokenizer = AutoTokenizer.from_pretrained(LLM_MODEL_PATH, mirror="modelscope")
            self.model = AutoModelForCausalLM.from_pretrained(LLM_MODEL_PATH, ms_dtype=mindspore.float32, mirror="modelscope")
            
            logging.info("✅ LLMActor: 模型加载成功!")
        except Exception as e:
            logging.info(f"❌ LLMActor: 加载LLM模型时发生严重错误: {e}")
            self.model = None
            self.tokenizer = None

    def generate(self, prompt: str, max_length: int = 1024) -> str:
        if not self.model or not self.tokenizer:
            return "LLM Actor 模型未成功加载，无法生成回答。"
            
        print(f"⚙️ LLMActor: 收到生成请求，正在调用 model.chat...")
        try:
            response_text, _ = self.model.chat(
                self.tokenizer,
                prompt,
                history=[],
                max_length=max_length,
            )
            return response_text
        except Exception as e:
            print(f"❌ LLMActor: 推理过程中发生错误: {e}")
            return f"LLM Actor 推理失败: {str(e)}"

# 3. parse_and_chunk_document Ray Task
@ray.remote
def parse_and_chunk_document(file_content: bytes, file_name: str) -> List[str]:
    if isinstance(file_content, str):
        logging.warning(f"⚠️ Ray Task: 文件 '{file_name}' 的内容是字符串(str)，正在自动编码为字节(bytes)。")
        file_content = file_content.encode('utf-8')
    print(f"⚙️ Ray Task: 正在解析文件 '{file_name}'...")
    from rapidocr_onnxruntime import RapidOCR
    from PIL import Image
    import numpy as np
    ocr_engine = RapidOCR()

    def read_pdf(stream) -> str:
        reader = PyPDF2.PdfReader(stream)
        return "".join(page.extract_text() for page in reader.pages if page.extract_text())

    def read_markdown(stream) -> str:
        md_text = stream.read().decode('utf-8')
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, 'html.parser')
        return re.sub(r'http\S+', '', soup.get_text())

    def read_text(stream) -> str:
        return stream.read().decode('utf-8')

    def read_image(stream) -> str:
        from rapidocr_onnxruntime import RapidOCR
        from PIL import Image
        import numpy as np
        ocr_engine = RapidOCR()
        img = Image.open(stream)
        img_np = np.array(img)
        result, _ = ocr_engine(img_np)
        if result:
            return "\n".join([line[1] for line in result])
        return ""

    def read_docx(stream) -> str:
        document = docx.Document(stream)
        return "\n".join([para.text for para in document.paragraphs])

    def read_pptx(stream) -> str:
        presentation = pptx.Presentation(stream)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs)

    def read_csv(stream) -> str:
        df = pd.read_csv(stream)
        return df.to_string()

    def read_audio(stream, file_name) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file_name)[1]) as tmp_file:
            tmp_file.write(stream.read())
            tmp_file_path = tmp_file.name
        
        try:
            print("🎤 Ray Task: 正在加载 Whisper 模型进行语音识别...")
            model = whisper.load_model("small")
            result = model.transcribe(tmp_file_path, fp16=False)
            print(f"🔊 Ray Task: 文件 '{file_name}' 语音识别完成。")
            return result.get("text", "")
        finally:
            os.remove(tmp_file_path) 


    file_stream = BytesIO(file_content)
    text = ""
    file_suffix = file_name.split('.')[-1].lower()

    if file_suffix == 'pdf': text = read_pdf(file_stream)
    elif file_suffix == 'md': text = read_markdown(file_stream)
    elif file_suffix == 'txt': text = read_text(file_stream)
    elif file_suffix in ['png', 'jpg', 'jpeg']: text = read_image(file_stream)
    elif file_suffix == 'docx': text = read_docx(file_stream)
    elif file_suffix == 'pptx': text = read_pptx(file_stream)
    elif file_suffix == 'csv': text = read_csv(file_stream)
    elif file_suffix in ['wav', 'mp3', 'm4a']: text = read_audio(file_stream, file_name)
    else:
        print(f"⚠️ Ray Task: 不支持的文件类型 '{file_suffix}'，跳过文件 {file_name}。")
        return []

    if not text:
        print(f"ℹ️ Ray Task: 从文件 '{file_name}' 中未提取到文本。")
        return []

    chunks = []

    if file_suffix == 'md':
        print(f"✨ Ray Task: 对 Markdown 文件 '{file_name}' 使用标题分割策略。")
        headers_to_split_on = [
            ("#", "Header 1"),
            ("##", "Header 2"),
            ("###", "Header 3"),
        ]
        markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on, strip_headers=False)
        md_header_splits = markdown_splitter.split_text(text)
        
        chunk_size = 600
        chunk_overlap = 150
        text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            encoding_name="cl100k_base",
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
        for split in md_header_splits:
            chunks.extend(text_splitter.split_text(split.page_content))

    else:
        print(f"✨ Ray Task: 对文件 '{file_name}' 使用递归字符分割策略。")
        chunk_size = 600
        chunk_overlap = 150
        text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            encoding_name="cl100k_base",
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
        chunks = text_splitter.split_text(text)
        
    final_chunks = [
        {
            "content": chunk_text,
            "source": file_name,
            "chunk_index": i + 1
        }
        for i, chunk_text in enumerate(chunks)
    ]
        
    print(f"✅ Ray Task: 文件 '{file_name}' 解析并分块为 {len(final_chunks)} 块。")
    return final_chunks